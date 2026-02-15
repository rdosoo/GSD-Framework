"""Generate the Get Shit Done Framework PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour Palette ──
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY = RGBColor(0x0F, 0x1A, 0x33)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GREY = RGBColor(0x99, 0x99, 0x99)
SOFT_WHITE = RGBColor(0xE8, 0xEE, 0xF4)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_accent_bar(slide, left, top, width, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_NAVY)

# Accent line
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(6), TEAL)

# Title
add_text_box(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.5),
             "GET SHIT DONE", font_size=54, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

# Subtitle
add_text_box(slide, Inches(2), Inches(3.5), Inches(9.333), Inches(1),
             "A Context Engineering Framework for Reliable AI-Augmented Development",
             font_size=22, color=TEAL, alignment=PP_ALIGN.CENTER)

# Author
add_text_box(slide, Inches(2), Inches(5.0), Inches(9.333), Inches(0.5),
             "github.com/glittercowboy/get-shit-done  |  MIT License  |  v1.11",
             font_size=14, color=MID_GREY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "GSD is an open-source meta-prompting and context engineering system that makes AI coding assistants reliable for building production software. Created by solo developer TACHES, trusted by engineers at Amazon, Google, Shopify, and Webflow.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
             "THE CONTEXT ROT PROBLEM", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(3))

# Quality zones - visual representation
zones = [
    ("0-30%", "PEAK", GREEN, Inches(1.5)),
    ("30-50%", "GOOD", YELLOW, Inches(1.2)),
    ("50-70%", "DEGRADING", ORANGE, Inches(0.9)),
    ("70%+", "POOR", RED, Inches(0.6)),
]

x_start = Inches(0.8)
y_start = Inches(2.0)
bar_width = Inches(2.5)

for i, (label, quality, color, height) in enumerate(zones):
    y = y_start + Inches(i * 1.15)
    bar = add_rect(slide, x_start, y, bar_width, height, color)
    add_text_box(slide, x_start + Inches(0.2), y + Pt(4), Inches(2), Inches(0.4),
                 f"{label}  {quality}", font_size=14, color=DARK_NAVY, bold=True)

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(4), Inches(0.5),
             "Context Window Usage \u2192 Quality Degrades", font_size=13, color=MID_GREY)

# Right side: problem bullets
add_bullet_list(slide, Inches(6.5), Inches(2.0), Inches(6), Inches(4.5), [
    "\u25b6  AI output degrades as context fills up",
    "\u25b6  \"Vibecoding\" produces inconsistent results at scale",
    "\u25b6  Existing tools add enterprise overhead",
    "\u25b6  No systematic approach to context management",
    "\u25b6  Session continuity is lost across resets",
], font_size=18, spacing=Pt(14))

add_notes(slide, "Claude Code is powerful, but quality degrades predictably with context usage. At 70%+ context, you get inconsistent, error-prone code. Existing spec-driven tools don't address this \u2014 they add sprint ceremonies and enterprise processes. GSD was built to solve context rot.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: What Is GSD?
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "WHAT IS GSD?", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(2))

# Left column - What You See
card1 = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
                  RGBColor(0x15, 0x22, 0x3E), TEAL)
add_text_box(slide, Inches(1.2), Inches(2.0), Inches(4.5), Inches(0.5),
             "WHAT YOU SEE", font_size=20, color=TEAL, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(2.7), Inches(4.5), Inches(3.5), [
    "npx get-shit-done-cc",
    "/gsd:new-project",
    "/gsd:plan-phase 1",
    "/gsd:execute-phase 1",
    "/gsd:verify-work 1",
], font_size=17, color=SOFT_WHITE, spacing=Pt(12))

# Right column - What's Happening
card2 = add_shape(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(5),
                  RGBColor(0x15, 0x22, 0x3E), ORANGE)
add_text_box(slide, Inches(7.4), Inches(2.0), Inches(4.5), Inches(0.5),
             "WHAT'S HAPPENING", font_size=20, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(2.7), Inches(4.5), Inches(3.5), [
    "11 specialised AI agents",
    "Context window management",
    "XML prompt formatting",
    "Subagent orchestration",
    "State persistence & recovery",
], font_size=17, color=SOFT_WHITE, spacing=Pt(12))

# Bottom tagline
add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.5),
             "Supports Claude Code  \u2022  OpenCode  \u2022  Gemini CLI   |   Mac, Windows, Linux",
             font_size=14, color=MID_GREY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "GSD hides significant complexity behind simple slash commands. You interact with a few commands; behind the scenes, the system manages 11 specialised agents, context windows, structured artifacts, and verification pipelines. Installs in one command, works cross-platform.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Core Philosophy
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "DESIGN PRINCIPLES", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(2.5))

principles = [
    ("PLANS = PROMPTS", "XML-structured executable instructions,\nnot documents to interpret", TEAL),
    ("FRESH CONTEXTS", "200k-token windows per execution\nunit \u2014 peak quality zone", GREEN),
    ("GOAL-BACKWARD", "Verify outcomes match intent,\nnot just task completion", ORANGE),
    ("AUTOMATION-FIRST", "Claude does the work;\nhumans verify and decide", TEAL),
    ("NO ENTERPRISE\nTHEATRE", "Built for builders, not\nbureaucrats", ORANGE),
]

for i, (title, desc, color) in enumerate(principles):
    x = Inches(0.6) + Inches(i * 2.5)
    y = Inches(2.0)

    card = add_shape(slide, x, y, Inches(2.3), Inches(4.5),
                     RGBColor(0x15, 0x22, 0x3E), color)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y + Inches(0.3),
                                     Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_NAVY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.15), y + Inches(1.2), Inches(2), Inches(0.9),
                 title, font_size=14, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(2.4), Inches(2), Inches(1.5),
                 desc, font_size=12, color=SOFT_WHITE,
                 alignment=PP_ALIGN.CENTER)

add_notes(slide, "Five principles drive GSD. Plans are literal prompts \u2014 XML-structured instructions. Every execution unit gets a fresh context. Verification checks goals, not tasks. Claude automates everything it can. No sprint ceremonies or story points.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: The Build Loop
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "THE CORE WORKFLOW", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(2.5))

# Init box at top
init_box = add_shape(slide, Inches(4.5), Inches(1.7), Inches(4.3), Inches(0.8), NAVY, TEAL)
tf = init_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "/gsd:new-project"
p.font.size = Pt(16)
p.font.color.rgb = TEAL
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Questions \u2192 Research \u2192 Requirements \u2192 Roadmap"
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GREY
p2.alignment = PP_ALIGN.CENTER

# Four workflow stages
stages = [
    ("DISCUSS", "/gsd:discuss-phase N", "Capture user\npreferences", TEAL),
    ("PLAN", "/gsd:plan-phase N", "Research + atomic\nplans + validation", GREEN),
    ("EXECUTE", "/gsd:execute-phase N", "Parallel waves\nfresh contexts", ORANGE),
    ("VERIFY", "/gsd:verify-work N", "Acceptance test\ngap analysis", RGBColor(0xAF, 0x7A, 0xC5)),
]

y_stage = Inches(3.2)
for i, (name, cmd, desc, color) in enumerate(stages):
    x = Inches(0.5) + Inches(i * 3.2)
    box = add_shape(slide, x, y_stage, Inches(2.8), Inches(2.5), RGBColor(0x15, 0x22, 0x3E), color)
    add_text_box(slide, x + Inches(0.2), y_stage + Inches(0.2), Inches(2.4), Inches(0.5),
                 name, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y_stage + Inches(0.8), Inches(2.4), Inches(0.4),
                 cmd, font_size=11, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y_stage + Inches(1.3), Inches(2.4), Inches(1),
                 desc, font_size=13, color=SOFT_WHITE, alignment=PP_ALIGN.CENTER)

    # Arrow between stages
    if i < 3:
        arrow_x = x + Inches(2.85)
        add_text_box(slide, arrow_x, y_stage + Inches(0.9), Inches(0.4), Inches(0.5),
                     "\u25b6", font_size=20, color=MID_GREY, alignment=PP_ALIGN.CENTER)

# Bottom: completion
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
             "Repeat per phase  \u2192  /gsd:complete-milestone  \u2192  /gsd:new-milestone  \u2192  Next cycle",
             font_size=15, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "The core workflow is a loop: initialise once, then discuss, plan, execute, verify per phase. Each phase gets user input, proper research, clean execution in fresh contexts, and goal-backward verification. When all phases complete, archive the milestone and start the next version.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Context Engineering
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "CONTEXT ENGINEERING", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(2.5))

# Artifact stack
artifacts = [
    ("PROJECT.md", "Project vision \u2014 always loaded", TEAL),
    ("REQUIREMENTS.md", "Scoped v1/v2 requirements", TEAL),
    ("ROADMAP.md", "Phased execution plan", TEAL),
    ("STATE.md", "Living memory across sessions", ORANGE),
    ("CONTEXT.md", "User\u2019s implementation decisions", GREEN),
    ("PLAN.md", "XML-structured executable prompts", GREEN),
    ("VERIFICATION.md", "Goal achievement report", RGBColor(0xAF, 0x7A, 0xC5)),
]

for i, (name, desc, color) in enumerate(artifacts):
    y = Inches(1.8) + Inches(i * 0.72)
    bar = add_shape(slide, Inches(0.8), y, Inches(5.5), Inches(0.6),
                    RGBColor(0x15, 0x22, 0x3E), color)
    add_text_box(slide, Inches(1.0), y + Pt(4), Inches(2.2), Inches(0.4),
                 name, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(3.2), y + Pt(4), Inches(3), Inches(0.4),
                 desc, font_size=13, color=SOFT_WHITE)

# Right side: key insight
insight_box = add_shape(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(5),
                        RGBColor(0x15, 0x22, 0x3E), TEAL)
add_text_box(slide, Inches(7.4), Inches(2.0), Inches(4.7), Inches(0.5),
             "KEY INSIGHT", font_size=18, color=TEAL, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(2.7), Inches(4.7), Inches(3.8), [
    "Every artifact has size constraints mapped to Claude's quality curve",
    "STATE.md is read FIRST in every workflow, updated after every action",
    "Plans stay at 2-3 tasks \u2014 small enough for the peak quality zone",
    "Each executor gets a fresh 200k-token context with zero accumulated garbage",
    "Result: consistent quality throughout the entire project lifecycle",
], font_size=14, color=SOFT_WHITE, spacing=Pt(14))

add_notes(slide, "GSD manages Claude's context through structured artifacts with deliberate size constraints. PROJECT.md holds the vision. STATE.md is living memory read first everywhere. Plans are small enough for peak quality. Each executor gets fresh 200k tokens. The result is consistent output quality.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Multi-Agent Architecture
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "MULTI-AGENT ARCHITECTURE", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(3))

# Central orchestrator
orch = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(3.0),
                               Inches(2.3), Inches(1.5))
orch.fill.solid()
orch.fill.fore_color.rgb = NAVY
orch.line.color.rgb = TEAL
orch.line.width = Pt(2)
tf = orch.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "THIN\nORCHESTRATOR"
p.font.size = Pt(14)
p.font.color.rgb = TEAL
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Agent quadrants
quadrants = [
    ("RESEARCH", ["Project Researcher", "Phase Researcher", "Synthesizer", "Codebase Mapper"],
     Inches(0.5), Inches(1.7), TEAL),
    ("PLANNING", ["Planner", "Plan Checker", "Roadmapper", ""],
     Inches(8.5), Inches(1.7), GREEN),
    ("EXECUTION", ["Executor (parallel)", "Fresh 200k context", "Per-task commits", ""],
     Inches(0.5), Inches(4.8), ORANGE),
    ("VERIFICATION", ["Verifier", "Integration Checker", "Debugger", ""],
     Inches(8.5), Inches(4.8), RGBColor(0xAF, 0x7A, 0xC5)),
]

for title, agents, x, y, color in quadrants:
    card = add_shape(slide, x, y, Inches(4.2), Inches(2.3),
                     RGBColor(0x15, 0x22, 0x3E), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.8), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    agent_items = [a for a in agents if a]
    add_bullet_list(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.8), Inches(1.5),
                    agent_items, font_size=12, color=SOFT_WHITE, spacing=Pt(4))

# Bottom stat
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
             "Main context stays at 30-40% usage \u2014 heavy lifting happens in subagent contexts",
             font_size=13, color=MID_GREY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "11 specialised agents coordinated by thin orchestrators. Research agents investigate in parallel. Planner creates plans, checker validates in a loop. Executors get fresh 200k-token contexts. Verifier confirms goals achieved. The orchestrator only spawns, waits, and integrates \u2014 your main session stays fast.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: XML Prompt Formatting
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "PLANS AS EXECUTABLE PROMPTS", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(3))

# Code block
code_bg = add_shape(slide, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8),
                    RGBColor(0x0A, 0x12, 0x28), TEAL)

code_lines = [
    '<task type="auto">',
    '  <name>Create login endpoint</name>',
    '  <files>src/api/auth/login.ts</files>',
    '  <action>',
    '    POST endpoint: {email, password}.',
    '    Query User by email, compare with',
    '    bcrypt. JWT via jose library.',
    '    Set httpOnly cookie on success.',
    '  </action>',
    '  <verify>',
    '    curl -X POST localhost:3000/api/auth',
    '    returns 200 + Set-Cookie header',
    '  </verify>',
    '  <done>Valid creds -> cookie, invalid -> 401</done>',
    '</task>',
]

code_text = "\n".join(code_lines)
code_box = add_text_box(slide, Inches(1.1), Inches(2.0), Inches(6), Inches(4.2),
                        code_text, font_size=13, color=SOFT_WHITE, font_name="Consolas")

# Right side annotations
annotations = [
    ("TASK TYPE", "auto | checkpoint:human-verify\ncheckpoint:decision", TEAL),
    ("FILES", "Exact targets \u2014\nno ambiguity", GREEN),
    ("ACTION", "Precise instructions with\nlibraries and approach", ORANGE),
    ("VERIFY", "Concrete test command\nbuilt into every task", RGBColor(0xAF, 0x7A, 0xC5)),
    ("DONE", "Measurable acceptance\ncriteria", TEAL),
]

for i, (label, desc, color) in enumerate(annotations):
    y = Inches(1.8) + Inches(i * 1.0)
    add_text_box(slide, Inches(8), y, Inches(2), Inches(0.3),
                 label, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(8), y + Inches(0.3), Inches(4.5), Inches(0.6),
                 desc, font_size=12, color=SOFT_WHITE)

# Bottom note
add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4),
             "2-3 tasks per plan \u2014 small enough for peak quality zone  |  Verification built into every task",
             font_size=13, color=MID_GREY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "Plans are structured XML optimised for Claude. Each task specifies exact files, precise actions, a verification command, and acceptance criteria. 2-3 tasks per plan keeps each executor in the peak quality zone. There's no ambiguity \u2014 Claude knows exactly what to build and how to verify it.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Wave-Based Execution
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "WAVE-BASED PARALLEL EXECUTION", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(3))

# Wave visualisation
waves = [
    ("WAVE 1", ["Plan A: Auth endpoints", "Plan B: Database schema"], TEAL, "parallel"),
    ("WAVE 2", ["Plan C: Frontend (depends on A+B)"], GREEN, "sequential"),
    ("WAVE 3", ["Plan D: Integration tests"], ORANGE, "sequential"),
]

for i, (wave_name, plans, color, mode) in enumerate(waves):
    y = Inches(1.9) + Inches(i * 1.6)
    # Wave label
    add_text_box(slide, Inches(0.8), y, Inches(1.5), Inches(0.4),
                 wave_name, font_size=16, color=color, bold=True)
    # Plan boxes
    for j, plan in enumerate(plans):
        x = Inches(2.8) + Inches(j * 3.5)
        box = add_shape(slide, x, y, Inches(3.2), Inches(0.7),
                        RGBColor(0x15, 0x22, 0x3E), color)
        add_text_box(slide, x + Inches(0.2), y + Pt(6), Inches(2.8), Inches(0.4),
                     plan, font_size=13, color=SOFT_WHITE)
        # Fresh context badge
        add_text_box(slide, x + Inches(0.2), y + Inches(0.4), Inches(2.8), Inches(0.3),
                     "\u26a1 Fresh 200k context", font_size=10, color=MID_GREY)

# Right side: git commits
git_box = add_shape(slide, Inches(7.5), Inches(1.9), Inches(5), Inches(3.5),
                    RGBColor(0x0A, 0x12, 0x28), GREEN)
add_text_box(slide, Inches(7.8), Inches(2.1), Inches(4), Inches(0.4),
             "ATOMIC GIT COMMITS", font_size=16, color=GREEN, bold=True)

commits = [
    "abc123 feat(01-01): create auth endpoints",
    "def456 feat(01-01): add database schema",
    "hij789 feat(01-02): build frontend forms",
    "lmn012 test(01-03): add integration tests",
]
add_bullet_list(slide, Inches(7.8), Inches(2.7), Inches(4.5), Inches(2.5),
                commits, font_size=12, color=SOFT_WHITE, spacing=Pt(10))

# Bottom benefits
benefits = [
    "Parallel within waves, sequential across waves",
    "Every task = one atomic commit",
    "Git bisect finds exact failing task",
    "Walk away, come back to completed work",
]
y_bottom = Inches(5.8)
for i, benefit in enumerate(benefits):
    x = Inches(0.5) + Inches(i * 3.2)
    add_shape(slide, x, y_bottom, Inches(3), Inches(0.9),
              RGBColor(0x15, 0x22, 0x3E), MID_GREY)
    add_text_box(slide, x + Inches(0.15), y_bottom + Pt(8), Inches(2.7), Inches(0.7),
                 benefit, font_size=12, color=SOFT_WHITE, alignment=PP_ALIGN.CENTER)

add_notes(slide, "Execution is wave-based. Independent plans run in parallel, each in a fresh 200k-token context. Every completed task gets an atomic git commit. Walk away, come back to completed work with a clean, bisectable git history. Main context stays light.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Verification Pipeline
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "GOAL-BACKWARD VERIFICATION", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(3))

# Three verification levels - pyramid style
levels = [
    ("LEVEL 3: WIRED", "Connected to the system\nComponent\u2192API, API\u2192DB, Form\u2192Handler",
     Inches(3.5), Inches(5.5), RGBColor(0xAF, 0x7A, 0xC5)),
    ("LEVEL 2: SUBSTANTIVE", "Real implementation, not stubs\nNo TODOs, placeholders, or hardcoded values",
     Inches(2.2), Inches(7.5), ORANGE),
    ("LEVEL 1: EXISTS", "File/component is present in the codebase",
     Inches(1), Inches(9.5), GREEN),
]

y_base = Inches(2.0)
for i, (title, desc, x, width, color) in enumerate(levels):
    y = y_base + Inches(i * 1.5)
    box = add_shape(slide, x, y, width, Inches(1.2), RGBColor(0x15, 0x22, 0x3E), color)
    add_text_box(slide, x + Inches(0.3), y + Pt(4), width - Inches(0.6), Inches(0.3),
                 title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.4), width - Inches(0.6), Inches(0.7),
                 desc, font_size=12, color=SOFT_WHITE, alignment=PP_ALIGN.CENTER)

# Right side: key principle
principle_box = add_shape(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(4.5),
                          RGBColor(0x15, 0x22, 0x3E), RED)
add_text_box(slide, Inches(8.8), Inches(2.2), Inches(3.6), Inches(0.5),
             "CORE PRINCIPLE", font_size=18, color=RED, bold=True)
add_bullet_list(slide, Inches(8.8), Inches(2.9), Inches(3.6), Inches(3.2), [
    "Task completion \u2260 Goal achievement",
    "Never trust SUMMARY claims blindly",
    "Detect stubs: TODOs, placeholders, empty returns",
    "Gaps feed back into planner for closure",
], font_size=13, color=SOFT_WHITE, spacing=Pt(14))

add_notes(slide, "Verification is goal-backward. Level 1 checks files exist. Level 2 confirms real implementations, not stubs. Level 3 verifies wiring between components. If gaps are found, they feed back into the planner. The system never declares success just because tasks ran.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: Configuration & Model Profiles
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "CONFIGURATION & MODEL PROFILES", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(3))

# Model profiles table
profiles = [
    ("PROFILE", "PLANNING", "EXECUTION", "VERIFICATION"),
    ("Quality", "Opus", "Opus", "Sonnet"),
    ("Balanced", "Opus", "Sonnet", "Sonnet"),
    ("Budget", "Sonnet", "Sonnet", "Haiku"),
]

table_top = Inches(1.8)
col_w = Inches(2.8)
row_h = Inches(0.55)

for r, row in enumerate(profiles):
    for c, cell in enumerate(row):
        x = Inches(0.8) + Inches(c * 2.8)
        y = table_top + Inches(r * 0.6)
        if r == 0:
            bg_color = TEAL
            text_color = DARK_NAVY
            bold = True
        else:
            bg_color = RGBColor(0x15, 0x22, 0x3E)
            text_color = SOFT_WHITE
            bold = (c == 0)
        cell_shape = add_rect(slide, x, y, col_w, row_h, bg_color)
        add_text_box(slide, x + Inches(0.15), y + Pt(4), col_w - Inches(0.3), row_h,
                     cell, font_size=14, color=text_color, bold=bold,
                     alignment=PP_ALIGN.CENTER)

# Right side: other settings
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.4),
             "OTHER SETTINGS", font_size=18, color=TEAL, bold=True)

settings = [
    ("Execution Mode", "interactive (confirm) | yolo (auto-approve)"),
    ("Planning Depth", "quick (1-3) | standard (3-5) | comprehensive (5-10)"),
    ("Workflow Agents", "Research, Plan Checker, Verifier \u2014 all toggleable"),
    ("Git Branching", "none | per-phase | per-milestone"),
    ("Parallelisation", "Enabled by default \u2014 independent plans run simultaneously"),
]

for i, (label, desc) in enumerate(settings):
    y = Inches(5.1) + Inches(i * 0.45)
    add_text_box(slide, Inches(0.8), y, Inches(2.5), Inches(0.4),
                 label, font_size=13, color=ORANGE, bold=True)
    add_text_box(slide, Inches(3.5), y, Inches(9), Inches(0.4),
                 desc, font_size=13, color=SOFT_WHITE)

add_notes(slide, "GSD is configurable. Model profiles balance quality vs API cost. Quality uses Opus everywhere, Balanced puts Opus on planning, Budget minimises Opus. You can also control planning depth, toggle workflow agents, set execution mode, and choose git branching strategy.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: Quick Mode & Session Management
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "FLEXIBILITY BUILT IN", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(2.5))

# Three feature panels
panels = [
    ("QUICK MODE", "/gsd:quick", [
        "Ad-hoc tasks with GSD guarantees",
        "Same planner + executor agents",
        "Skips research & verification",
        "Bug fixes, small features, config",
    ], TEAL, "\u26a1"),
    ("SESSION MANAGEMENT", "/gsd:pause-work\n/gsd:resume-work", [
        "Pause mid-phase, resume later",
        "Full context restoration",
        "STATE.md tracks everything",
        "No lost work across resets",
    ], ORANGE, "\u23ef"),
    ("BROWNFIELD SUPPORT", "/gsd:map-codebase", [
        "Analyse existing codebases",
        "Parallel agents: stack, arch, conventions",
        "new-project knows your patterns",
        "Questions focus on what's new",
    ], GREEN, "[B]"),
]

for i, (title, cmd, items, color, icon) in enumerate(panels):
    x = Inches(0.5) + Inches(i * 4.2)
    card = add_shape(slide, x, Inches(1.7), Inches(3.9), Inches(5.3),
                     RGBColor(0x15, 0x22, 0x3E), color)
    add_text_box(slide, x + Inches(0.3), Inches(1.9), Inches(3.3), Inches(0.4),
                 f"{icon}  {title}", font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(0.5),
                 cmd, font_size=12, color=LIGHT_GREY)
    add_bullet_list(slide, x + Inches(0.3), Inches(3.3), Inches(3.3), Inches(3.2),
                    items, font_size=13, color=SOFT_WHITE, spacing=Pt(10))

add_notes(slide, "Not every task needs full planning. Quick mode provides atomic commits and state tracking for small tasks. Session management handles pause/resume across context resets. For existing codebases, map-codebase analyses your stack so the system knows your patterns.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: Impact & Adoption
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "WHY THIS MATTERS", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(2.5))

# Value cards
cards = [
    ("CONSISTENT\nQUALITY", "Fresh contexts prevent\ndegradation. Verification\nensures goals are met.", TEAL),
    ("FULL\nTRACEABILITY", "Atomic commits. Structured\nartifacts. Every decision\ndocumented.", GREEN),
    ("MULTI-RUNTIME\nSUPPORT", "Claude Code, OpenCode,\nGemini CLI. No vendor\nlock-in.", ORANGE),
    ("OPEN SOURCE\nMIT LICENSE", "Active community.\nFast evolution. Used\nat top tech companies.", RGBColor(0xAF, 0x7A, 0xC5)),
]

for i, (title, desc, color) in enumerate(cards):
    x = Inches(0.5) + Inches(i * 3.2)
    y = Inches(1.8)
    card = add_shape(slide, x, y, Inches(3), Inches(3), RGBColor(0x15, 0x22, 0x3E), color)
    # Top accent bar
    add_rect(slide, x, y, Inches(3), Pt(4), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.4), Inches(2.4), Inches(0.8),
                 title, font_size=18, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.5), Inches(2.4), Inches(1.3),
                 desc, font_size=13, color=SOFT_WHITE,
                 alignment=PP_ALIGN.CENTER)

# Bottom: audience-specific benefits
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.4),
             "FOR TECHNICAL LEADERSHIP", font_size=16, color=TEAL, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.5), [
    "\u25b6  Predictable, verifiable output from AI coding assistants \u2014 not random vibecoding",
    "\u25b6  Configurable quality/cost tradeoff via model profiles \u2014 practical for any budget",
    "\u25b6  Reproducible process \u2014 same commands, same workflow, consistent results across developers",
], font_size=14, color=SOFT_WHITE, spacing=Pt(8))

add_notes(slide, "For technical leadership: GSD provides predictable quality through verification, full traceability through atomic commits and structured artifacts, multi-runtime support without vendor lock-in, and it's MIT licensed with an active community. It's a practical tool for organisations exploring AI-augmented development.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: Getting Started
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "GET STARTED", font_size=36, color=WHITE, bold=True,
             font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(2))

# Steps
steps = [
    ("1", "INSTALL", "npx get-shit-done-cc", TEAL),
    ("2", "VERIFY", "/gsd:help", GREEN),
    ("3", "INITIALISE", "/gsd:new-project  (or /gsd:map-codebase first)", ORANGE),
    ("4", "BUILD", "discuss \u2192 plan \u2192 execute \u2192 verify", RGBColor(0xAF, 0x7A, 0xC5)),
    ("5", "SHIP", "/gsd:complete-milestone", TEAL),
]

for i, (num, label, cmd, color) in enumerate(steps):
    y = Inches(1.7) + Inches(i * 0.95)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Pt(4),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_NAVY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.6), y + Pt(2), Inches(1.5), Inches(0.5),
                 label, font_size=18, color=color, bold=True)
    # Command in monospace box
    cmd_bg = add_shape(slide, Inches(3.3), y, Inches(6), Inches(0.6),
                       RGBColor(0x0A, 0x12, 0x28), RGBColor(0x33, 0x44, 0x66))
    add_text_box(slide, Inches(3.5), y + Pt(4), Inches(5.6), Inches(0.4),
                 cmd, font_size=14, color=SOFT_WHITE, font_name="Consolas")

# Resources section
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.4),
             "RESOURCES", font_size=16, color=TEAL, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4),
             "GitHub: github.com/glittercowboy/get-shit-done   |   NPM: get-shit-done-cc   |   Discord: discord.gg/5JJgD5svVS",
             font_size=14, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "Getting started is one command: npx get-shit-done-cc. Verify with /gsd:help, start your first project, follow the core loop, and ship. The GitHub repo has full documentation and the Discord community is active. MIT licensed and evolving fast.")


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = r"C:\Users\dosoor\Projects\GSD-Framework\get-shit-done-framework.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
